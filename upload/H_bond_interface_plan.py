# -*- coding: utf-8 -*-
# H_bond_interface.py — ChimeraX 1.9+
# 1) HBonds：全部（>18 行自动 6 列）+ 按链对分页（空则不出页），并同时写出 SVG
# 2) Salt bridges：同 HBonds（saltOnly true），同样分页和 SVG
# 3) Interface：每个链对两页（表格 或 文本框）+ measured/fixed SVG；SVG 单独落盘；
#    如装了 cairosvg，会把 SVG 转 PNG 插入 PPT

from chimerax.core.commands import run as cxrun
from chimerax.atomic import selected_atoms, selected_residues, AtomicStructure
from chimerax.core.colors import color_name
import os, re, tempfile, math, traceback, itertools
from collections import Counter

# =============== 可调参数 ===============
INCLUDE_NUCLEIC_NUMBER = True
ROW_THRESHOLD_FOR_6COL = 18
IFACE_COLOR_TEXT   = True     # 接口残基（表格/文本框）是否用链色，False = 全黑
IFACE_TEXTBOX_MODE = True     # True=每个残基单独文本框；False=表格
SVG_COLS = 30
SVG_UNMODELED_ALPHA = 0.35
SVG_MEASURED_USE_CHAIN_COLORS = True   # measured.svg 是否用链色
SVG_FIXED_USE_CHAIN_COLORS    = False  # fixed.svg 是否用链色

DEFAULT_FASTA = "/Users/lixingyue/Desktop/pengbo/iggend/sequence.fa"  # 自动回退（不会删除或改动你的 FASTA）

# =============== 通用工具 ===============
def q(p):
    p = os.path.abspath(os.path.expanduser(str(p)))
    return f'"{p}"' if any(ch in p for ch in ' ()') else p

def _get_kv(args, key, default=None):
    for i in range(len(args) - 1):
        if str(args[i]).lower() == key:
            return os.path.expanduser(args[i + 1])
    return default

def _has_flag(args, key):
    return any(str(a).lower() == key for a in args)

def _derive_out_paths(out_val, script_dir):
    if out_val:
        out_abs = os.path.abspath(os.path.expanduser(out_val))
        if out_abs.lower().endswith(".pptx"):
            ppt_path = out_abs
            base_dir = os.path.dirname(out_abs)
        else:
            base_dir = out_abs; os.makedirs(base_dir, exist_ok=True)
            ppt_path = os.path.join(base_dir, "hbonds_table.pptx")
    else:
        base_dir = script_dir
        ppt_path = os.path.join(base_dir, "hbonds_table.pptx")
    save_txt_path = os.path.join(base_dir, "hbond.txt")
    svg_dir = os.path.join(base_dir, "pair_svgs"); os.makedirs(svg_dir, exist_ok=True)
    return ppt_path, save_txt_path, svg_dir

def _chain_id_from_label(label: str) -> str:
    try: return str(label).split("_", 1)[0]
    except Exception: return "?"

def base_mid(mid):
    if not mid: return ""
    return "#" + mid.lstrip("#").split(".")[0]

# =============== 颜色映射 ===============
NAME2HEX = {
    "mediumslateblue":"#7B68EE","lightcoral":"#F08080","darkseagreen":"#8FBC8F",
    "red":"#FF0000","green":"#008000","blue":"#0000FF","cyan":"#00FFFF","magenta":"#FF00FF",
    "yellow":"#FFFF00","orange":"#FFA500","gold":"#FFD700","dodgerblue":"#1E90FF","deepskyblue":"#00BFFF",
    "royalblue":"#4169E1","slateblue":"#6A5ACD","mediumpurple":"#9370DB","purple":"#800080","plum":"#DDA0DD",
    "orchid":"#DA70D6","lime":"#00FF00","limegreen":"#32CD32","seagreen":"#2E8B57","forestgreen":"#228B22",
    "darkgreen":"#006400","springgreen":"#00FF7F","turquoise":"#40E0D0","teal":"#008080","cadetblue":"#5F9EA0",
    "steelblue":"#4682B4","lightskyblue":"#87CEFA","navy":"#000080","maroon":"#800000","brown":"#A52A2A",
    "salmon":"#FA8072","darksalmon":"#E9967A","tomato":"#FF6347","coral":"#FF7F50","lightpink":"#FFB6C1",
    "hotpink":"#FF69B4","pink":"#FFC0CB","indianred":"#CD5C5C","crimson":"#DC143C","chocolate":"#D2691E",
    "sienna":"#A0522D","peru":"#CD853F","khaki":"#F0E68C","tan":"#D2B48C","wheat":"#F5DEB3",
    "olive":"#808000","darkolivegreen":"#556B2F","olivedrab":"#6B8E23","chartreuse":"#7FFF00","greenyellow":"#ADFF2F",
    "lightgreen":"#90EE90","palegreen":"#98FB98","mediumseagreen":"#3CB371","mediumspringgreen":"#00FA9A",
    "aquamarine":"#7FFFD4","lightsalmon":"#FFA07A","peachpuff":"#FFDAB9","bisque":"#FFE4C4","lavender":"#E6E6FA",
    "thistle":"#D8BFD8","slategray":"#708090","lightgray":"#D3D3D3","gray":"#808080","black":"#000000","white":"#FFFFFF"
}
PALETTE = ["#1F77B4","#FF7F0E","#2CA02C","#D62728","#9467BD","#8C564B","#E377C2","#7F7F7F","#BCBD22","#17BECF"]

def _normalize_css_color(s):
    if not s: return ("#111111", 1.0)
    s = s.strip().lower()
    if s.startswith("#"):
        if len(s) == 9: return ("#" + s[1:7], int(s[7:9], 16) / 255.0)
        if len(s) == 7: return (s, 1.0)
    return (NAME2HEX.get(s, "#111111"), 1.0)

def _palette_by_cid(cid):
    try: idx = (ord(str(cid)[0]) if cid else 0) % len(PALETTE)
    except Exception: idx = 0
    return PALETTE[idx]

def build_chain_color_map(session, logger=None):
    cmap_by_model, cmap_by_chain = {}, {}
    models = session.models.list(type=AtomicStructure)
    for m in models:
        mid_full = m.id_string
        mid_base = base_mid(mid_full)
        for mid in (mid_full, mid_base): cmap_by_model.setdefault(mid, {})
        chains = sorted(set(r.chain_id for r in m.residues))
        for cid in chains:
            samples, cnt = [], 0
            for r in m.residues:
                if r.chain_id != cid: continue
                c = None
                try:
                    if r.ribbon_color is not None: c = r.ribbon_color
                except Exception: pass
                if c is None:
                    try:
                        if len(r.atoms) > 0: c = r.atoms[0].color
                    except Exception: pass
                if c is not None:
                    try: samples.append(color_name(c))
                    except Exception: pass
                cnt += 1
                if cnt >= 400: break
            if samples:
                raw = Counter(samples).most_common(1)[0][0]
            else:
                try: raw = color_name(m.color)
                except Exception: raw = None
            hx, a = _normalize_css_color(raw) if raw else ("#111111", 1.0)
            if hx == "#111111" or raw is None:
                hx = _palette_by_cid(cid); a = 1.0
                if logger: logger.warning(f"[chain-color] fallback -> {mid_base}/{cid} = {hx}")
            for mid in (mid_full, mid_base):
                cmap_by_model[mid][cid] = {"raw": raw or "None", "hex": hx, "alpha": a}
            cmap_by_chain[cid] = {"hex": hx, "alpha": a}
            if logger: logger.info(f"[chain-color] {mid_base}/{cid} -> {hx}")
    return cmap_by_model, cmap_by_chain

def pick_color(mid, cid, cmap_by_model, cmap_by_chain):
    for key in (mid, base_mid(mid)):
        if key in cmap_by_model and cid in cmap_by_model[key]:
            return cmap_by_model[key][cid]["hex"], 1.0
    if cid in cmap_by_chain:
        return cmap_by_chain[cid]["hex"], 1.0
    return _palette_by_cid(cid), 1.0

def _hex_to_rgb_tuple(h):
    h = h.lstrip("#")
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# =============== AA 3→1 & label ============
AA3_TO_1 = {
    "ALA":"A","ARG":"R","ASN":"N","ASP":"D","CYS":"C","GLU":"E","GLN":"Q","GLY":"G",
    "HIS":"H","ILE":"I","LEU":"L","LYS":"K","MET":"M","PHE":"F","PRO":"P","SER":"S",
    "THR":"T","TRP":"W","TYR":"Y","VAL":"V","SEC":"U","PYL":"O","MSE":"M"
}

def _format_res_label(resn: str, resi: str) -> str:
    rn = (resn or "").upper()
    if rn in AA3_TO_1: return f"{AA3_TO_1[rn]}{resi}"
    if len(rn)==2 and rn[0]=="D" and rn[1] in "ACGTU":
        return f"DNA/{rn[1]}{resi if INCLUDE_NUCLEIC_NUMBER else ''}"
    if rn in {"A","C","G","U","T"}:
        return f"RNA/{rn}{resi if INCLUDE_NUCLEIC_NUMBER else ''}"
    return f"{rn}{resi}"

# =============== HBond 日志解析（含回退） ===============
HBOND_ROW = re.compile(
    r"""/([A-Za-z])\s+([A-Z0-9]{1,3})\s*(-?\d+)\s+([A-Za-z0-9']+)\s+
         /([A-Za-z])\s+([A-Z0-9]{1,3})\s*(-?\d+)\s+([A-Za-z0-9']+)\s+
         no\shydrogen\s+([0-9.]+)
     """, re.VERBOSE
)

def parse_hbond_log(text: str):
    def _fmt2(x: str):
        try: return f"{float(x):.2f}"
        except Exception: return x
    rows=[]
    for m in HBOND_ROW.finditer(text):
        d_chain,d_resn,d_resi,d_atom,a_chain,a_resn,a_resi,a_atom,dist = m.groups()
        donor    = f"{d_chain}_{_format_res_label(d_resn, d_resi)}_{d_atom}"
        acceptor = f"{a_chain}_{_format_res_label(a_resn, a_resi)}_{a_atom}"
        rows.append((donor, acceptor, _fmt2(dist)))
    return rows

def parse_hbond_log_fallback(text: str):
    rows=[]
    for line in text.splitlines():
        s=line.strip()
        if not s.startswith("/"): continue
        parts=s.split()
        if len(parts)<10: continue
        try:
            d_chain=parts[0][1:]; d_resn, d_resi, d_atom = parts[1], parts[2], parts[3]
            a_chain=parts[4][1:]; a_resn, a_resi, a_atom = parts[5], parts[6], parts[7]
            dist_tok = parts[-2] if parts[-1].upper()=="N/A" else parts[-1]
            dist = f"{float(dist_tok):.2f}"
            donor    = f"{d_chain}_{_format_res_label(d_resn, d_resi)}_{d_atom}"
            acceptor = f"{a_chain}_{_format_res_label(a_resn, a_resi)}_{a_atom}"
            rows.append((donor, acceptor, dist))
        except Exception:
            continue
    return rows

# =============== HBond 表 → SVG（简洁风） ===============
def _write_hbond_svg(rows, svg_path, *, sixcol_threshold=ROW_THRESHOLD_FOR_6COL):
    os.makedirs(os.path.dirname(svg_path) or ".", exist_ok=True)
    single_block = (len(rows) <= sixcol_threshold)

    W = 1200; H0 = 80
    line_gap = 24
    margin_l, margin_t = 40, 40

    def block_cols(x0, width):
        d = int(x0 + width * 0.00)
        a = int(x0 + width * 0.44)
        dist = int(x0 + width * 0.88)
        return d, a, dist

    if single_block:
        rows_per_col = len(rows)
        height = H0 + rows_per_col * line_gap + 40
        with open(svg_path, "w", encoding="utf-8") as f:
            f.write(f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{height}" viewBox="0 0 {W} {height}">\n')
            f.write('<style>text{font-family:"Times New Roman",serif;fill:#000}</style>\n')
            x0 = margin_l; y = margin_t
            c1, c2, c3 = block_cols(x0, W - 2*margin_l)
            def head(x, main):
                f.write(f'<text x="{x}" y="{y}" font-size="16">{main}</text>\n')
                f.write(f'<text x="{x}" y="{y+16}" font-size="12"><tspan baseline-shift="sub">(chain_residue_atom)</tspan></text>\n')
            head(c1, "donor"); head(c2, "acceptor")
            f.write(f'<text x="{c3}" y="{y}" font-size="16">distance (Å)</text>\n')
            y = y + 36
            for d, a, dist in rows:
                f.write(f'<text x="{c1}" y="{y}" font-size="14">{d}</text>\n')
                f.write(f'<text x="{c2}" y="{y}" font-size="14">{a}</text>\n')
                f.write(f'<text x="{c3}" y="{y}" font-size="14">{dist}</text>\n')
                y += line_gap
            f.write('</svg>\n')
    else:
        half = (W - 2*margin_l)
        left_w = half // 2
        rows_left = math.ceil(len(rows)/2)
        rows_right = len(rows) - rows_left
        rows_per_col = max(rows_left, rows_right)
        height = H0 + rows_per_col * line_gap + 40
        with open(svg_path, "w", encoding="utf-8") as f:
            f.write(f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{height}" viewBox="0 0 {W} {height}">\n')
            f.write('<style>text{font-family:"Times New Roman",serif;fill:#000}</style>\n')
            y = margin_t
            c1, c2, c3 = block_cols(margin_l, left_w)
            c4, c5, c6 = block_cols(margin_l + left_w + 20, left_w - 20)
            def head(x, main):
                f.write(f'<text x="{x}" y="{y}" font-size="16">{main}</text>\n')
                f.write(f'<text x="{x}" y="{y+16}" font-size="12"><tspan baseline-shift="sub">(chain_residue_atom)</tspan></text>\n')
            head(c1, "donor"); head(c2, "acceptor"); f.write(f'<text x="{c3}" y="{y}" font-size="16">distance (Å)</text>\n')
            head(c4, "donor"); head(c5, "acceptor"); f.write(f'<text x="{c6}" y="{y}" font-size="16">distance (Å)</text>\n')
            y = y + 36
            for i in range(rows_per_col):
                if i < rows_left:
                    d,a,dist = rows[i]
                    f.write(f'<text x="{c1}" y="{y}" font-size="14">{d}</text>\n')
                    f.write(f'<text x="{c2}" y="{y}" font-size="14">{a}</text>\n')
                    f.write(f'<text x="{c3}" y="{y}" font-size="14">{dist}</text>\n')
                j = i + rows_left
                if j < len(rows):
                    d2,a2,dist2 = rows[j]
                    f.write(f'<text x="{c4}" y="{y}" font-size="14">{d2}</text>\n')
                    f.write(f'<text x="{c5}" y="{y}" font-size="14">{a2}</text>\n')
                    f.write(f'<text x="{c6}" y="{y}" font-size="14">{dist2}</text>\n')
                y += line_gap
            f.write('</svg>\n')

# =============== PPT: HBond 表（无填充、括号下标）==============
def _ppt_common_styles():
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    return RGBColor, Inches, Pt

def darken_hex_color(hex_str, factor=0.8):
    hex_str = hex_str.lstrip("#")
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    r = max(0, min(255, int(r * factor)))
    g = max(0, min(255, int(g * factor)))
    b = max(0, min(255, int(b * factor)))
    return f"#{r:02X}{g:02X}{b:02X}"

def _add_hbond_table_slide(prs, rows, *, caption=None):
    """
    固定 3 列（donor / acceptor / distance），把整张表放在 PPT 的右侧，
    给左侧留出充足空间放图片。自动按内容分配三列宽度。
    """
    RGBColor, Inches, Pt = _ppt_common_styles()
    from pptx.enum.text import PP_ALIGN

    # === 新建空白页 & 白背景 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === 顶部小标题（左对齐，可选）===
    if caption:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.4))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = caption
        r.font.name = "Times New Roman"
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

    # === 右侧区域：把表格整体放到右边 ===
    # 常见 10x7.5 英寸页面上，留出 ~4.8 英寸给左侧图片
    left, top, width, height = Inches(5.0), Inches(1.0), Inches(4.8), Inches(6.0)

    # --------- 内部小工具 ----------
    def est_width_in(chars, avg_char_in=0.082, padding_in=0.26):
        # 估计一列需要的宽度（英寸）
        return max(padding_in, chars * avg_char_in + padding_in)

    def compute_col_widths(total_in, rows, min_triplet, max_triplet):
        # 根据内容长度自适应三列宽度（但始终只有 3 列）
        max_d = max((len(d) for d, _, _ in rows), default=len("donor"))
        max_a = max((len(a) for _, a, _ in rows), default=len("acceptor"))
        max_dist = max((len(s) for *_, s in rows), default=len("distance (Å)"))
        w_dist = est_width_in(max(max_dist, len("distance (Å)")))
        w_dist = min(max_triplet[2], max(min_triplet[2], w_dist))

        rem = max(0.1, total_in - w_dist)
        w_d = est_width_in(max_d)
        w_a = est_width_in(max_a)
        scale = rem / max(w_d + w_a, 0.01)
        w_d = min(max_triplet[0], max(min_triplet[0], w_d * scale))
        w_a = min(max_triplet[1], max(min_triplet[1], w_a * scale))

        used = w_d + w_a + w_dist
        if abs(used - total_in) > 1e-3:
            diff = total_in - used
            w_d += diff / 2
            w_a += diff / 2
        return (w_d, w_a, w_dist)

    BLACK = RGBColor(0, 0, 0)

    def style_cell(cell, text, size_pt=12, bold=False, nowrap=False):
        cell.text = str(text)
        cell.fill.background()
        tf = cell.text_frame
        tf.word_wrap = not nowrap
        if not tf.paragraphs:
            tf.add_paragraph()
        for p in tf.paragraphs:
            if not p.runs:
                p.add_run()
            for r in p.runs:
                r.font.name = "Times New Roman"
                r.font.size = Pt(size_pt)
                r.font.bold = bool(bold)
                r.font.color.rgb = BLACK

    def set_header_with_sub(cell, main_text, sub_text, *, size_pt=12, nowrap=False):
        cell.fill.background()
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = not nowrap
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = main_text
        r1.font.name = "Times New Roman"
        r1.font.size = Pt(size_pt)
        r1.font.color.rgb = BLACK

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = sub_text
        r2.font.name = "Times New Roman"
        r2.font.size = Pt(size_pt)
        r2.font.color.rgb = BLACK
        r2.font.subscript = True

    def neutralize(tbl):
        for a in ("first_row", "last_row", "first_col", "last_col", "banded_rows", "banded_columns"):
            if hasattr(tbl, a):
                setattr(tbl, a, False)

    # === 只有 3 列，不再根据行数改为 6 列 ===
    n = len(rows)
    tbl = slide.shapes.add_table(n + 1, 3, left, top, width, height).table
    neutralize(tbl)

    # 列宽：限制区间可以按需要微调
    min_triplet = (1.6, 1.6, 0.9)
    max_triplet = (3.2, 3.2, 1.4)
    w_d, w_a, w_dist = compute_col_widths(width.inches, rows, min_triplet, max_triplet)
    tbl.columns[0].width = Inches(w_d)
    tbl.columns[1].width = Inches(w_a)
    tbl.columns[2].width = Inches(w_dist)

    # 行高 & 字号：简单自适应，尽量容纳更多行
    header_h_in = 0.28
    row_h_in = 0.22
    font_pt = 12
    if n > 28:
        row_h_in = 0.20
        font_pt = 11
    if n > 36:
        row_h_in = 0.18
        font_pt = 10

    tbl.rows[0].height = Inches(header_h_in)
    for r in range(1, len(tbl.rows)):
        tbl.rows[r].height = Inches(row_h_in)

    # 表头
    set_header_with_sub(tbl.cell(0, 0), "donor", "(chain_residue_atom)")
    set_header_with_sub(tbl.cell(0, 1), "acceptor", "(chain_residue_atom)")
    style_cell(tbl.cell(0, 2), "distance (Å)", nowrap=True, size_pt=font_pt)

    # 统一清除填充
    for row in tbl.rows:
        for cell in row.cells:
            cell.fill.background()

    # 数据
    for i, (d, a, dist) in enumerate(rows, start=1):
        style_cell(tbl.cell(i, 0), d, size_pt=font_pt)
        style_cell(tbl.cell(i, 1), a, size_pt=font_pt)
        style_cell(tbl.cell(i, 2), dist, size_pt=font_pt, nowrap=True)

def add_hbond_all_and_pairs(prs, rows, svg_dir):
    if not rows: return
    _add_hbond_table_slide(prs, rows, caption=None)
    _write_hbond_svg(rows, os.path.join(svg_dir, "hbonds_all.svg"))
    chains = sorted({ _chain_id_from_label(d) for d,_,_ in rows } | { _chain_id_from_label(a) for _,a,_ in rows })
    for c1, c2 in itertools.combinations(chains, 2):
        subset = [r for r in rows if { _chain_id_from_label(r[0]), _chain_id_from_label(r[1]) } == {c1, c2}]
        if subset:
            _add_hbond_table_slide(prs, subset, caption=f"HBonds — Chains {c1}–{c2}")
            _write_hbond_svg(subset, os.path.join(svg_dir, f"hbonds_{c1}-{c2}.svg"))

# =============== Salt bridges：与 HBonds 同构的入口（表格+SVG） ===============
def add_salt_all_and_pairs(prs, rows, svg_dir):
    if not rows: return
    _add_hbond_table_slide(prs, rows, caption="Salt bridges — all")
    _write_hbond_svg(rows, os.path.join(svg_dir, "salts_all.svg"))
    chains = sorted({ _chain_id_from_label(d) for d,_,_ in rows } | { _chain_id_from_label(a) for _,a,_ in rows })
    for c1, c2 in itertools.combinations(chains, 2):
        subset = [r for r in rows if { _chain_id_from_label(r[0]), _chain_id_from_label(r[1]) } == {c1, c2}]
        if subset:
            _add_hbond_table_slide(prs, subset, caption=f"Salt bridges — Chains {c1}–{c2}")
            _write_hbond_svg(subset, os.path.join(svg_dir, f"salts_{c1}-{c2}.svg"))

# =============== Interface：残基获取（选中接口） ===============
def _iface_residues_for_pair(session, c1, c2):
    def _pick(chain):
        out=[]
        sel = selected_residues(session)
        if sel:
            for r in sel:
                if getattr(r, "chain_id", None) == chain:
                    out.append((getattr(r,"name","UNK"), str(getattr(r,"number","?"))))
        seen=set(); uniq=[]
        for nm,num in sorted(out, key=lambda x: int(re.sub(r"[^-0-9]","",x[1]) or 0)):
            key=(nm,num)
            if key not in seen:
                uniq.append((nm,num)); seen.add(key)
        return uniq
    try: cxrun(session, "select clear")
    except Exception: pass
    session.logger.info(f"RUN: interfaces select /{c1} contacting /{c2} bothSides true")
    cxrun(session, f"interfaces select /{c1} contacting /{c2} bothSides true")
    r1 = _pick(c1); r2 = _pick(c2)
    try: cxrun(session, "select clear")
    except Exception: pass
    return r1, r2

def _iface_label(nm3, num):
    nm3u=(nm3 or "").upper()
    if nm3u in AA3_TO_1: return f"{AA3_TO_1[nm3u]}{num}"
    if len(nm3u)==2 and nm3u[0]=="D" and nm3u[1] in "ACGTU":
        return f"DNA/{nm3u[1]}{num if INCLUDE_NUCLEIC_NUMBER else ''}"
    if nm3u in {"A","C","G","U","T"}:
        return f"RNA/{nm3u}{num if INCLUDE_NUCLEIC_NUMBER else ''}"
    return f"{nm3u}{num}"

# =============== Interface：表格页（无填充） ===============
def _add_interface_table_slide(prs, session, cmap_by_model, cmap_by_chain, c1, c2, left_labels, right_labels):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255,255,255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.4))
    tf = tb.text_frame; tf.clear()
    r = tf.paragraphs[0].add_run(); r.text = f"Interface residues — Chains {c1}–{c2}"
    r.font.name="Times New Roman"; r.font.size=Pt(12); r.font.color.rgb = RGBColor(0,0,0)

    left, top, width, height = Inches(0.5), Inches(1.0), Inches(8.5), Inches(5.9)
    n_rows = max(len(left_labels), len(right_labels)) + 1
    tbl = slide.shapes.add_table(n_rows, 2, left, top, width, height).table

    for a in ("first_row","last_row","first_col","last_col","banded_rows","banded_columns"):
        if hasattr(tbl,a): setattr(tbl,a,False)

    tbl.columns[0].width = Inches(4.25); tbl.columns[1].width = Inches(4.25)
    for c,h in enumerate((f"{c1} residues", f"{c2} residues")):
        cell = tbl.cell(0,c); cell.text=h; cell.fill.background()
        for p in cell.text_frame.paragraphs:
            if not p.runs: p.add_run()
            for rr in p.runs:
                rr.font.name="Times New Roman"; rr.font.size=Pt(12); rr.font.bold=True; rr.font.color.rgb=RGBColor(0,0,0)

    models = session.models.list(type=AtomicStructure)
    first_mid = models[0].id_string if models else "#1"
    hex1,_ = pick_color(first_mid, c1, cmap_by_model, cmap_by_chain)
    hex2,_ = pick_color(first_mid, c2, cmap_by_model, cmap_by_chain)
    def hex_to_rgb(h): h=h.lstrip("#"); return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    col1 = RGBColor(*hex_to_rgb(hex1)); col2 = RGBColor(*hex_to_rgb(hex2))

    def set_cell(ridx, cidx, text, color_rgb):
        cell = tbl.cell(ridx, cidx); cell.text=str(text); cell.fill.background()
        for p in cell.text_frame.paragraphs:
            if not p.runs: p.add_run()
            for rr in p.runs:
                rr.font.name="Times New Roman"; rr.font.size=Pt(12)
                rr.font.color.rgb = (color_rgb if IFACE_COLOR_TEXT else RGBColor(0,0,0))

    for i in range(1, n_rows):
        set_cell(i,0, left_labels[i-1]  if i-1<len(left_labels)  else "", col1)
        set_cell(i,1, right_labels[i-1] if i-1<len(right_labels) else "", col2)

    for row in tbl.rows:
        for cell in row.cells:
            cell.fill.background()

# =============== Interface：文本框版（每个残基独立文本框 + 自动分页） ===============
def _add_interface_labels_slide(prs, c1, c2, left_labels, right_labels,
                                left_rgb=None, right_rgb=None):
    """
    单页右侧布局的文本框版接口残基：
    - 右侧区域显示两条链（左链=left_labels，右链=right_labels）
    - 每条链在自己的半宽区域内，最多使用两列（同一页内）
    - 如果一列放不下，会自动起第二列；如两列仍放不下，会自动减小行距/字号以塞进两列
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ---- 画布与区域（右侧） ----
    PAGE_TOP_IN      = 1.0
    PAGE_TITLE_TOP   = 0.5
    RIGHT_LEFT_IN    = 5.0     # 整体右移：右侧区域起点（留出左侧放图片）
    AREA_W_IN        = 4.8     # 右侧区域宽度
    AREA_H_IN        = 6.0     # 右侧区域高度
    CHAIN_GAP_IN     = 0.30    # 两条链（左右半区）之间的间隙
    SUBCOL_GAP_IN    = 0.20    # 同一链内两个子列之间的间隔

    # 文本外观
    FONT_NAME        = "Times New Roman"
    FONT_PT_BASE     = 12
    LINE_H_BASE_IN   = 0.26
    FONT_PT_MIN      = 10
    LINE_H_MIN_IN    = 0.18

    # 颜色
    BLACK = RGBColor(0, 0, 0)
    left_color  = RGBColor(*(left_rgb  or (0, 0, 0)))
    right_color = RGBColor(*(right_rgb or (0, 0, 0)))

    # 右侧区域分成左右两半（各自一条链）
    half_w = (AREA_W_IN - CHAIN_GAP_IN) / 2.0
    left_area_x  = RIGHT_LEFT_IN
    right_area_x = RIGHT_LEFT_IN + half_w + CHAIN_GAP_IN

    # 估计在给定行高下，单列可容纳的行数
    def max_rows_per_col(line_h_in):
        return max(1, int((AREA_H_IN - (PAGE_TOP_IN - PAGE_TOP_IN)) / line_h_in))

    # 根据当前行高/字号，决定是否需要两列；若两列仍放不下则逐步压缩
    def fit_two_columns(labels):
        font_pt = FONT_PT_BASE
        line_h  = LINE_H_BASE_IN

        while True:
            m = max_rows_per_col(line_h)
            cols_needed = (len(labels) + m - 1) // m
            if cols_needed <= 2:
                return m, 2 if len(labels) > m else 1, font_pt, line_h
            # 压缩：优先减小行高，其次减小字号
            if line_h > LINE_H_MIN_IN + 1e-6:
                line_h = max(LINE_H_MIN_IN, line_h - 0.02)
            elif font_pt > FONT_PT_MIN:
                font_pt -= 1
            else:
                # 已到最小仍放不下，也强行两列（最后几条可能略超区但尽量容纳）
                return m, 2, font_pt, line_h

    # 将列表切成列块（每列最多 m 行，最多两列）
    def to_columns(labels, m):
        cols = []
        for i in range(0, min(len(labels), 2*m), m):
            cols.append(labels[i:i+m])
            if len(cols) == 2:
                break
        return cols

    # 绘制一条链（在给定区域 x0..x0+half_w）
    def draw_chain(slide, x0_in, chain_name, labels, rgb_color):
        # 先拟合两列范围
        m, ncols, font_pt, line_h = fit_two_columns(labels)
        cols = to_columns(labels, m)

        # 标题（链名）
        tb = slide.shapes.add_textbox(Inches(x0_in), Inches(PAGE_TOP_IN - 0.30),
                                      Inches(half_w), Inches(0.25))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{chain_name} residues"
        r.font.name = FONT_NAME
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT

        # 列宽：按 ncols=1 或 2 分配
        if ncols == 1:
            sub_w = half_w
            xs = [x0_in]
        else:
            sub_w = (half_w - SUBCOL_GAP_IN) / 2.0
            xs = [x0_in, x0_in + sub_w + SUBCOL_GAP_IN]

        # 逐列绘制
        for col_idx, col_labels in enumerate(cols):
            x = xs[col_idx]
            y = PAGE_TOP_IN
            for t in col_labels:
                tb = slide.shapes.add_textbox(Inches(x), Inches(y),
                                              Inches(sub_w), Inches(line_h))
                tf = tb.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = str(t)
                r.font.name = FONT_NAME
                r.font.size = Pt(font_pt)
                r.font.color.rgb = (rgb_color if IFACE_COLOR_TEXT else BLACK)
                y += line_h

    # === 新建空白页 & 白背景 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 顶部标题
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(PAGE_TITLE_TOP), Inches(9.0), Inches(0.4))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Interface residues — Chains {c1}–{c2} (textbox mode)"
    r.font.name = FONT_NAME
    r.font.size = Pt(12)
    r.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT

    # 绘制两条链（都在右侧区域，且各自最多两列）
    draw_chain(slide, left_area_x,  c1, left_labels,  left_color)
    draw_chain(slide, right_area_x, c2, right_labels, right_color)

# =============== Interface SVG ===============

# 解析 FASTA 头：比如 E[-62--1]
HDR_RE = re.compile(r'^(?P<name>[^\s\[\]]+)(?:\[(?P<start>-?\d+)\s*-\s*(?P<end>-?\d+)\])?$')

def read_fasta(path):
    """
    返回 dict: { name: {"seq": SEQ, "start": start_int, "end": end_int} }
    - 无范围时，start=1, end=len(seq)
    - 有范围时，严格按 header 给的 start/end；如果与序列长度不符，就把 end 纠正为 start+len(seq)-1
    """
    if not path or not os.path.exists(path): return {}
    out, cur_name, cur_start, cur_end, parts = {}, None, None, None, []
    with open(path, "r", encoding="utf-8") as f:
        for line in f:
            s=line.strip()
            if not s:
                continue
            if s.startswith(">"):
                if cur_name is not None:
                    seq = "".join(parts).replace(" ","").replace("\t","").upper()
                    if cur_start is None: cur_start = 1
                    if cur_end   is None: cur_end   = cur_start + len(seq) - 1
                    if (cur_end - cur_start + 1) != len(seq):
                        cur_end = cur_start + len(seq) - 1
                    out[cur_name] = {"seq": seq, "start": int(cur_start), "end": int(cur_end)}
                parts=[]
                m = HDR_RE.match(s[1:].split()[0])
                if m:
                    cur_name = m.group("name")
                    cur_start = int(m.group("start")) if m.group("start") else None
                    cur_end   = int(m.group("end"))   if m.group("end")   else None
                else:
                    cur_name = s[1:].split()[0]
                    cur_start = None; cur_end = None
            else:
                parts.append(s)
    if cur_name is not None:
        seq = "".join(parts).replace(" ","").replace("\t","").upper()
        if cur_start is None: cur_start = 1
        if cur_end   is None: cur_end   = cur_start + len(seq) - 1
        if (cur_end - cur_start + 1) != len(seq):
            cur_end = cur_start + len(seq) - 1
        out[cur_name] = {"seq": seq, "start": int(cur_start), "end": int(cur_end)}
    return out

def resnums_to_seq_idx(nums, start, L):
    """
    nums: 结构残基号集合（可能有负数）
    start: FASTA header 解析出的起始编号（整数，可能为负）
    L: 序列长度
    输出: 1..L 内的索引集合
    """
    out=set()
    for n in nums:
        i = n - start + 1  # 映射到 1..L
        if 1 <= i <= L:
            out.add(i)
    return out

# ---- 重要：defattr 行解析用的正则（这在你原文缺失了）----
LINE_NUM = re.compile(
    r"#(?P<model>\d+)"
    r"/(?P<chain>[A-Za-z0-9])"
    r":(?P<num>-?\d+)"
)

def parse_defattr_list(path):
    out = {}
    if not os.path.exists(path): return out
    with open(path, "r", encoding="utf-8") as f:
        for line in f:
            m = LINE_NUM.search(line)
            if not m: continue
            mid_full = "#" + m.group("model").lstrip("#")
            mid = base_mid(mid_full)
            cid = m.group("chain"); num = int(m.group("num"))
            out.setdefault(mid, {}).setdefault(cid, set()).add(num)
    for mid in out:
        for cid in out[mid]:
            out[mid][cid] = sorted(out[mid][cid])
    return out

def _join(outdir, name): return os.path.abspath(os.path.join(outdir or ".", name))

def save_all_modeled(session, outdir):
    path = _join(outdir, "all.defattr")
    cxrun(session, f"save {q(path)} attrName r:number selectedOnly false matchMode any modelIds true")
    return path, parse_defattr_list(path)

def save_pair_iface(session, outdir, mid, c1, c2):
    try: cxrun(session, "select clear")
    except Exception: pass
    try:
        cxrun(session, f"interfaces select {mid}/{c1} contacting {mid}/{c2} bothSides true")
    except Exception:
        cxrun(session, f"interfaces select /{c1} contacting /{c2} bothSides true")
    path = _join(outdir, f"pair_{c1}_{c2}.defattr")
    cxrun(session, f"save {q(path)} attrName r:number selectedOnly true matchMode any modelIds true")
    return path, parse_defattr_list(path)

# ---- 统一且可用的 write_pair_svg（支持负号刻度）----
def write_pair_svg(svg_path, chain1, seq1, modeled1, iface1,
                   chain2, seq2, modeled2, iface2,
                   chunk=60, dx=12, line_gap=42,
                   color1="#111111", color1_alpha=1.0,
                   color2="#111111", color2_alpha=1.0,
                   unmodeled_alpha=0.35,
                   disp_start1=1, disp_start2=1):
    ACCENT= "#444"; left, right, top, sep = 80, 40, 46, 50
    modeled1=set(i for i in modeled1 if 1<=i<=len(seq1)); modeled2=set(i for i in modeled2 if 1<=i<=len(seq2))
    iface1=set(i for i in iface1 if 1<=i<=len(seq1));     iface2=set(i for i in iface2 if 1<=i<=len(seq2))
    lines1=max(1,(len(seq1)+chunk-1)//chunk); lines2=max(1,(len(seq2)+chunk-1)//chunk)
    width  = left + dx*chunk + right; height = top + lines1*line_gap + sep + lines2*line_gap + 40

    def draw_chain(f, label, seq, modeled, iface, y0, col_hex, alpha, disp_start):
        f.write(f'<text x="14" y="{y0}">{label}</text>\n')
        L=len(seq); rows=max(1,(L+chunk-1)//chunk); ua=max(0.0,min(1.0,alpha*unmodeled_alpha))
        for ln in range(rows):
            s=ln*chunk+1; e=min(L,(ln+1)*chunk); x0=left; y_seq=y0+ln*line_gap
            p=s
            while p<=e:
                x=x0+(p-s)*dx
                abs_num = disp_start + (p-1)
                f.write(f'<text x="{x}" y="{y_seq-14}" fill="{ACCENT}" font-size="12">{abs_num}</text>\n')
                p+=10
            for i,pos in enumerate(range(s,e+1)):
                x=x0+i*dx; aa=seq[pos-1]; op=alpha if pos in modeled else ua
                f.write(f'<text x="{x}" y="{y_seq}" fill="{col_hex}" style="fill:{col_hex}" fill-opacity="{op:.3f}">{aa}</text>\n')
                if pos in iface:
                    f.write(f'<line x1="{x}" y1="{y_seq+3}" x2="{x+dx-2}" y2="{y_seq+3}" stroke="{col_hex}" stroke-opacity="{alpha:.3f}" stroke-width="2.6"/>\n')

    os.makedirs(os.path.dirname(svg_path) or ".", exist_ok=True)
    with open(svg_path, "w", encoding="utf-8") as f:
        f.write(f'<!-- COLORS: {chain1}={color1} (a={color1_alpha:.2f}); {chain2}={color2} (a={color2_alpha:.2f}) -->\n')
        f.write(f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">\n')
        f.write('<style>text{font-family:Courier New,monospace;font-size:14px}</style>\n')
        draw_chain(f, chain1, seq1, set(modeled1), set(iface1), top, color1, color1_alpha, disp_start1)
        y2 = top + lines1*line_gap + sep
        draw_chain(f, chain2, seq2, set(modeled2), set(iface2), y2, color2, color2_alpha, disp_start2)
        f.write('</svg>\n')

def svg_to_png(svg_path):
    png_path = os.path.splitext(svg_path)[0] + ".png"
    try:
        import cairosvg
        cairosvg.svg2png(url=svg_path, write_to=png_path)
        return png_path
    except Exception:
        return None

def _add_interface_svg_slide(prs, measured_png, fixed_png, c1, c2):
    if not (measured_png or fixed_png): return
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255,255,255)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.4))
    tf = tb.text_frame; tf.clear()
    r = tf.paragraphs[0].add_run(); r.text = f"Interface diagrams — Chains {c1}–{c2}"
    r.font.name="Times New Roman"; r.font.size=Pt(12); r.font.color.rgb = RGBColor(0,0,0)
    y = Inches(1.0)
    if measured_png:
        slide.shapes.add_picture(measured_png, Inches(0.5), y, width=Inches(8.5)); y = y + Inches(3.1)
    if fixed_png:
        slide.shapes.add_picture(fixed_png, Inches(0.5), y, width=Inches(8.5))

# =============== Interface：总调度 ===============
def add_interface_pairs(prs, session, fasta_path, svg_dir):
    logger = getattr(session, "logger", None)
    cmap_by_model, cmap_by_chain = build_chain_color_map(session, logger)

    chains=set()
    for m in session.models.list(type=AtomicStructure):
        for r in m.residues:
            if r.chain_id: chains.add(r.chain_id)
    chains = sorted(chains)
    if not chains: return

    fasta = read_fasta(fasta_path) if fasta_path else {}

    models = session.models.list(type=AtomicStructure)
    mid_base = base_mid(models[0].id_string) if models else "#1"

    try:
        _, all_map = save_all_modeled(session, svg_dir)
    except Exception:
        all_map = {}

    for c1, c2 in itertools.combinations(chains, 2):
        r1_raw, r2_raw = _iface_residues_for_pair(session, c1, c2)
        left_labels  = [_iface_label(nm, num) for nm,num in r1_raw]
        right_labels = [_iface_label(nm, num) for nm,num in r2_raw]

        if left_labels or right_labels:
            if IFACE_TEXTBOX_MODE:
                hex1,_ = pick_color(mid_base, c1, cmap_by_model, cmap_by_chain)
                hex2,_ = pick_color(mid_base, c2, cmap_by_model, cmap_by_chain)
                _add_interface_labels_slide(
                    prs, c1, c2, left_labels, right_labels,
                    _hex_to_rgb_tuple(hex1), _hex_to_rgb_tuple(hex2)
                )
            else:
                _add_interface_table_slide(prs, session, cmap_by_model, cmap_by_chain, c1, c2, left_labels, right_labels)
        else:
            session.logger.info(f"Skip empty interface page for {c1}-{c2}")

        measured_png = fixed_png = None
        try:
            try:
                _, pair_map = save_pair_iface(session, svg_dir, mid_base, c1, c2)
            except Exception:
                pair_map = {}

            fi1 = fasta.get(c1) or {}
            fi2 = fasta.get(c2) or {}
            seq1 = fi1["seq"] if isinstance(fi1, dict) else (fi1 or "")
            seq2 = fi2["seq"] if isinstance(fi2, dict) else (fi2 or "")

            start1 = (fi1.get("start") if isinstance(fi1, dict) else 1) or 1
            start2 = (fi2.get("start") if isinstance(fi2, dict) else 1) or 1

            modeled1_raw = set(all_map.get(mid_base, {}).get(c1, []))
            modeled2_raw = set(all_map.get(mid_base, {}).get(c2, []))
            iface1_raw   = set(pair_map.get(mid_base, {}).get(c1, []))
            iface2_raw   = set(pair_map.get(mid_base, {}).get(c2, []))

            modeled1 = resnums_to_seq_idx(modeled1_raw, start1, len(seq1))
            modeled2 = resnums_to_seq_idx(modeled2_raw, start2, len(seq2))
            iface1   = resnums_to_seq_idx(iface1_raw,   start1, len(seq1))
            iface2   = resnums_to_seq_idx(iface2_raw,   start2, len(seq2))

            if seq1 and seq2:
                # --- measured（可选链色）---
                if SVG_MEASURED_USE_CHAIN_COLORS:
                    m_hex1, _ = pick_color(mid_base, c1, cmap_by_model, cmap_by_chain)
                    m_hex2, _ = pick_color(mid_base, c2, cmap_by_model, cmap_by_chain)
                else:
                    m_hex1, _ = _normalize_css_color("black")
                    m_hex2, _ = _normalize_css_color("black")

                svg_meas = os.path.join(svg_dir, f"pair_{c1}_{c2}_measured.svg")
                write_pair_svg(
                    svg_meas, c1, seq1, modeled1, iface1,
                    c2, seq2, modeled2, iface2,
                    chunk=SVG_COLS, color1=m_hex1, color1_alpha=1.0,
                    color2=m_hex2, color2_alpha=1.0,
                    unmodeled_alpha=SVG_UNMODELED_ALPHA,
                    disp_start1=start1, disp_start2=start2
                )
                measured_png = svg_to_png(svg_meas)

                # --- fixed（默认黑白或链色，按配置）---
                if SVG_FIXED_USE_CHAIN_COLORS:
                    f_hex1, _ = pick_color(mid_base, c1, cmap_by_model, cmap_by_chain)
                    f_hex2, _ = pick_color(mid_base, c2, cmap_by_model, cmap_by_chain)
                else:
                    f_hex1, _ = _normalize_css_color("black")
                    f_hex2, _ = _normalize_css_color("black")

                svg_fix = os.path.join(svg_dir, f"pair_{c1}_{c2}_fixed.svg")
                write_pair_svg(
                    svg_fix, c1, seq1, modeled1, iface1,
                    c2, seq2, modeled2, iface2,
                    chunk=SVG_COLS, color1=f_hex1, color1_alpha=1.0,
                    color2=f_hex2, color2_alpha=1.0,
                    unmodeled_alpha=SVG_UNMODELED_ALPHA,
                    disp_start1=start1, disp_start2=start2
                )
                fixed_png = svg_to_png(svg_fix)
            else:
                session.logger.info(f"No FASTA for chains {c1}/{c2}; SVG skipped")
        except Exception as e:
            session.logger.info(f"SVG stage failed for {c1}-{c2}: {e}")

        if measured_png or fixed_png:
            _add_interface_svg_slide(prs, measured_png, fixed_png, c1, c2)

# =============== 额外：按链对导出 H-bond/Salt 的视图 PNG ===============
def generate_hbond_salt_images(session, chain1, chain2, svg_dir, is_salt=False,
                               cmap_by_model=None, cmap_by_chain=None, mid=None):
    import os
    from chimerax.atomic import selected_atoms

    def save_pair(no_ext, w=2000, h=1675):
        cxrun(session, f"save {q(no_ext)}.cxs")
        cxrun(session, f"save {q(no_ext)}.png width {w} height {h} supersample 3 transparentBackground true")

    try:
        os.makedirs(svg_dir, exist_ok=True)

        if cmap_by_model is None or cmap_by_chain is None:
            _logger = getattr(session, "logger", None)
            cmap_by_model, cmap_by_chain = build_chain_color_map(session, _logger)

        if mid is None:
            models = session.models.list(type=AtomicStructure)
            mid = base_mid(models[0].id_string) if models else "#1"

        hex1, _ = pick_color(mid, chain1, cmap_by_model, cmap_by_chain)
        hex2, _ = pick_color(mid, chain2, cmap_by_model, cmap_by_chain)
        hex1 = darken_hex_color(hex1, 0.8)
        hex2 = darken_hex_color(hex2, 0.8)

        prefix = "salt" if is_salt else "hbond"
        base_path        = os.path.join(svg_dir, f"{prefix}_{chain1}_{chain2}")
        label_path       = f"{base_path}_label"        # cartoon + label
        atoms_label_path = f"{base_path}_atoms_label"  # atoms + label
        atoms_path       = f"{base_path}_atoms"        # atoms only

        # 统一外观
        cxrun(session, "color bychain")
        cxrun(session, "transparency 30 target r")
        cxrun(session, "set bgColor white")
        cxrun(session, "hide atoms; hide cartoons; hide surfaces; ~label")

        # 计算
        hb_base = f"hbonds /{chain1} restrict /{chain2} select true reveal true dashes 6 log true color black radius .05"
        cxrun(session, hb_base + (" saltOnly true" if is_salt else ""))

        # 1) Cartoon + label（先别清选择/别 ~label）
        #cxrun(session, f"show /{chain1} /{chain2} cartoons; style stick; color sel byhetero;view sel")
        #if len(selected_atoms(session)) == 0:
        #    cxrun(session, f"show /{chain1} /{chain2} cartoons")
        cxrun(session, f'label (sel & /{chain1}) text "{{0.label_one_letter_code}} {{0.number}}{{0.insertion_code}}" '
                       f'position primary offset 0,0,0 color {hex1}')
        cxrun(session, f'label (sel & /{chain2}) text "{{0.label_one_letter_code}} {{0.number}}{{0.insertion_code}}" '
                       f'position primary offset 0,0,0 color {hex2}')
        cxrun(session, "~key; select clear")                
        #save_pair(label_path)


        # 2) atoms + label
        cxrun(session, hb_base + (" saltOnly true" if is_salt else ""))
        cxrun(session, f"show /{chain1} /{chain2} cartoons; style stick; color sel byhetero;view sel;show sel atoms")
        if len(selected_atoms(session)) == 0:
            cxrun(session, f"show /{chain1} /{chain2} cartoons")
        else:
            cxrun(session, "show sel atoms;")
        cxrun(session, "~key; select clear")
        save_pair(atoms_label_path)


        # 3) atoms（无 label）
        cxrun(session, hb_base + (" saltOnly true" if is_salt else ""))
        if len(selected_atoms(session)) == 0:
            cxrun(session, f"~label")
        else:
            cxrun(session, "show sel atoms; ~label")
        cxrun(session, "~key; select clear")
        save_pair(atoms_path)
  

        return (f"Generated {prefix} images for {chain1}-{chain2}: "
                f"{os.path.basename(label_path)}.png, "
                f"{os.path.basename(atoms_label_path)}.png, "
                f"{os.path.basename(atoms_path)}.png")
    except Exception as e:
        #try: session.logger.error(f"[{prefix} {chain1}-{chain2}] export failed: {e}")
        #except Exception: pass
        return None


def add_hbond_salt_images_to_ppt(prs, session, chains, svg_dir, is_salt=False, rows=None):
    """
    每个链对：左侧叠加 PNG（atoms + atoms_label + cartoon_label），右侧三列表格（若提供 rows）。
    rows: 传入对应的 hbond/salt 解析结果（列表，每行: (donor, acceptor, dist)）
    """
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    import os

    cmap_by_model, cmap_by_chain = build_chain_color_map(session, getattr(session, "logger", None))
    models = session.models.list(type=AtomicStructure)
    mid = base_mid(models[0].id_string) if models else "#1"

    prefix = "salt" if is_salt else "hbond"
    title_prefix = "Salt" if is_salt else "H-bond"

    # 左侧图片区域
    IMG_LEFT_IN  = 0.5
    IMG_TOP_IN   = 2.0
    IMG_W_IN     = 4.8     # 左半屏宽度

    for c1, c2 in itertools.combinations(chains, 2):
        _ = generate_hbond_salt_images(session, c1, c2, svg_dir, is_salt,
                                       cmap_by_model=cmap_by_model,
                                       cmap_by_chain=cmap_by_chain,
                                       mid=mid)

        p_cartoon_label = os.path.join(svg_dir, f"{prefix}_{c1}_{c2}_label.png")
        p_atoms_label   = os.path.join(svg_dir, f"{prefix}_{c1}_{c2}_atoms_label.png")
        p_atoms         = os.path.join(svg_dir, f"{prefix}_{c1}_{c2}_atoms.png")

        if not any(os.path.exists(p) for p in (p_cartoon_label, p_atoms_label, p_atoms)):
            continue

        # 新建一页
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{title_prefix} {c1}-{c2}"
        for p in slide.shapes.title.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT

        # 左侧叠加 PNG
        if os.path.exists(p_atoms):
            slide.shapes.add_picture(p_atoms, Inches(IMG_LEFT_IN), Inches(IMG_TOP_IN), width=Inches(IMG_W_IN))
        if os.path.exists(p_atoms_label):
            slide.shapes.add_picture(p_atoms_label, Inches(IMG_LEFT_IN), Inches(IMG_TOP_IN), width=Inches(IMG_W_IN))
        #if os.path.exists(p_cartoon_label):
        #    slide.shapes.add_picture(p_cartoon_label, Inches(IMG_LEFT_IN), Inches(IMG_TOP_IN), width=Inches(IMG_W_IN))

        # 右侧三列表格（如果给了 rows，就筛这一对）
        if rows:
            pair_rows = [r for r in rows if { _chain_id_from_label(r[0]), _chain_id_from_label(r[1]) } == {c1, c2}]
            if pair_rows:
                _draw_hbond_table_on_slide(slide, pair_rows,
                                           left_in=5.0, top_in=1.0, width_in=4.8, height_in=6.0)


# =============== 额外：按链对导出 interface 的视图 PNG ===============
def generate_interface_images(session, chain1, chain2, svg_dir,
                               cmap_by_model=None, cmap_by_chain=None, mid=None):
    import os
    from chimerax.atomic import selected_atoms, selected_residues

    def cx(cmd):
        session.logger.info("RUN: " + cmd)
        try:
            cxrun(session, cmd)
        except Exception as e:
            session.logger.error(f"CMD FAIL: {cmd}\n -> {e}")
            raise

    def sel_stat(tag):
        aa = len(selected_atoms(session))
        rr = len(selected_residues(session))
        session.logger.info(f"[{tag}] selection -> atoms={aa}, residues={rr}")
        return aa, rr

    def save_png(path_no_ext, w=2000, h=1675):
        png_path = path_no_ext + ".png"
        cx(f"save {q(png_path)} width {w} height {h} supersample 3 transparentBackground true")
        if not os.path.exists(png_path):
            session.logger.error(f"[SAVE-FAIL] PNG not written: {png_path}")
        else:
            session.logger.info(f"[SAVE-OK] {png_path}")
        return png_path

    def save_cxs(path_no_ext):
        cxs_path = path_no_ext + ".cxs"
        cx(f"save {q(cxs_path)}")
        if not os.path.exists(cxs_path):
            session.logger.error(f"[SAVE-FAIL] CXS not written: {cxs_path}")
        else:
            session.logger.info(f"[SAVE-OK] {cxs_path}")
        return cxs_path

    try:
        os.makedirs(svg_dir, exist_ok=True)

        if cmap_by_model is None or cmap_by_chain is None:
            _logger = getattr(session, "logger", None)
            cmap_by_model, cmap_by_chain = build_chain_color_map(session, _logger)

        if mid is None:
            models = session.models.list(type=AtomicStructure)
            mid = base_mid(models[0].id_string) if models else "#1"

        hex1, _ = pick_color(mid, chain1, cmap_by_model, cmap_by_chain)
        hex2, _ = pick_color(mid, chain2, cmap_by_model, cmap_by_chain)
        hex1 = darken_hex_color(hex1, 0.8)
        hex2 = darken_hex_color(hex2, 0.8)

        base_path        = os.path.join(svg_dir, f"interface_{chain1}_{chain2}")
        label_path       = f"{base_path}_label"
        atoms_label_path = f"{base_path}_atoms_label"
        atoms_path       = f"{base_path}_atoms"

        cx("color bychain")
        cx("transparency 30 target r")
        cx("set bgColor white")
        cx("cartoon suppressBackboneDisplay 1; hide atoms; hide cartoons; hide surfaces; ~label")

        cx(f"interfaces select /{chain1} contacting /{chain2} bothSides true")
        aa, rr = sel_stat("iface-3")
        if aa == 0 and rr == 0:
            session.logger.warning(f"[iface-3] NO interface selection for {chain1}-{chain2} (atoms only)")
            cx(f"show /{chain1} /{chain2} cartoons")
        else:
            cx("show sel atoms; view sel; ~label; color sel byhetero")
        cx(f"show /{chain1} /{chain2} cartoons")
        cx("~key; select clear")

        save_cxs(atoms_path)
        cx("~key; select clear")
        save_png(atoms_path)

        #cx(f"interfaces select /{chain1} contacting /{chain2} bothSides true")
        #aa, rr = sel_stat("iface-1")
        #if aa == 0 and rr == 0:
        #    session.logger.warning(f"[iface-1] NO interface selection for {chain1}-{chain2}")
        #    cx(f"show /{chain1} /{chain2} cartoons")
        #else:
        #    cx(f"show /{chain1} /{chain2} cartoons")
        #cx("style stick; color sel byhetero;hide atoms")
        #cx(f'label (sel & /{chain1}) text "{{0.label_one_letter_code}} {{0.number}}{{0.insertion_code}}" position primary offset 0,0,0 color {hex1}')
        #cx(f'label (sel & /{chain2}) text "{{0.label_one_letter_code}} {{0.number}}{{0.insertion_code}}" position primary offset 0,0,0 color {hex2}')
        #save_cxs(label_path)
        #cx("hide atoms; ~key; select clear")
        #save_png(label_path)

        cx(f"interfaces select /{chain1} contacting /{chain2} bothSides true")
        cx(f'label (sel & /{chain1}) text "{{0.label_one_letter_code}} {{0.number}}{{0.insertion_code}}" position primary offset 0,0,0 color {hex1}')
        cx(f'label (sel & /{chain2}) text "{{0.label_one_letter_code}} {{0.number}}{{0.insertion_code}}" position primary offset 0,0,0 color {hex2}')
        cx(f"show sel atoms")
        aa, rr = sel_stat("iface-2")
        if aa == 0 and rr == 0:
            session.logger.warning(f"[iface-2] NO interface selection for {chain1}-{chain2} (atoms+label)")
            #cx(f"show /{chain1} /{chain2} atoms")
        else:
            cx("show sel atoms;color sel byhetero")
        save_cxs(atoms_label_path)
        cx("~key; select clear")
        save_png(atoms_label_path)

        return (f"Generated interface images for {chain1}-{chain2}: "
                f"{os.path.basename(label_path)}.png, "
                f"{os.path.basename(atoms_label_path)}.png, "
                f"{os.path.basename(atoms_path)}.png")

    except Exception:
        return None
def _draw_hbond_table_on_slide(slide, rows, *, left_in=5.0, top_in=1.0, width_in=4.8, height_in=6.0):
    """在给定 slide 的右侧区域绘制固定三列表格（donor / acceptor / distance）。"""
    RGBColor, Inches, Pt = _ppt_common_styles()

    # 小工具：自适应列宽（但始终三列）
    def est_width_in(chars, avg_char_in=0.082, padding_in=0.26):
        return max(padding_in, chars * avg_char_in + padding_in)

    def compute_col_widths(total_in, rows, min_triplet, max_triplet):
        max_d = max((len(d) for d, _, _ in rows), default=len("donor"))
        max_a = max((len(a) for _, a, _ in rows), default=len("acceptor"))
        max_dist = max((len(s) for *_, s in rows), default=len("distance (Å)"))
        w_dist = est_width_in(max(max_dist, len("distance (Å)")))
        w_dist = min(max_triplet[2], max(min_triplet[2], w_dist))
        rem = max(0.1, total_in - w_dist)
        w_d = est_width_in(max_d); w_a = est_width_in(max_a)
        scale = rem / max(w_d + w_a, 0.01)
        w_d = min(max_triplet[0], max(min_triplet[0], w_d * scale))
        w_a = min(max_triplet[1], max(min_triplet[1], w_a * scale))
        used = w_d + w_a + w_dist
        if abs(used - total_in) > 1e-3:
            diff = total_in - used
            w_d += diff / 2; w_a += diff / 2
        return (w_d, w_a, w_dist)

    def neutralize(tbl):
        for a in ("first_row","last_row","first_col","last_col","banded_rows","banded_columns"):
            if hasattr(tbl, a): setattr(tbl, a, False)

    BLACK = RGBColor(0,0,0)
    def style_cell(cell, text, size_pt=12, bold=False, nowrap=False):
        from pptx.util import Pt as _Pt
        cell.text = str(text); cell.fill.background()
        tf = cell.text_frame; tf.word_wrap = not nowrap
        if not tf.paragraphs: tf.add_paragraph()
        for p in tf.paragraphs:
            if not p.runs: p.add_run()
            for r in p.runs:
                r.font.name="Times New Roman"; r.font.size=_Pt(size_pt)
                r.font.bold=bool(bold); r.font.color.rgb = BLACK

    def set_header_with_sub(cell, main_text, sub_text, *, size_pt=12, nowrap=False):
        from pptx.util import Pt as _Pt
        cell.fill.background()
        tf = cell.text_frame; tf.clear(); tf.word_wrap = not nowrap
        p1 = tf.paragraphs[0]; r1=p1.add_run(); r1.text=main_text
        r1.font.name="Times New Roman"; r1.font.size=_Pt(size_pt); r1.font.color.rgb=BLACK
        p2 = tf.add_paragraph(); r2=p2.add_run(); r2.text=sub_text
        r2.font.name="Times New Roman"; r2.font.size=_Pt(size_pt); r2.font.color.rgb=BLACK; r2.font.subscript=True

    # 创建表格（固定三列）
    n = len(rows)
    left, top, width, height = Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    tbl = slide.shapes.add_table(n + 1, 3, left, top, width, height).table
    neutralize(tbl)

    # 列宽
    min_triplet=(1.6, 1.6, 0.9); max_triplet=(3.2, 3.2, 1.4)
    w_d,w_a,w_dist = compute_col_widths(width.inches, rows, min_triplet, max_triplet)
    tbl.columns[0].width = Inches(w_d); tbl.columns[1].width = Inches(w_a); tbl.columns[2].width = Inches(w_dist)

    # 行高/字号（多行时更紧凑）
    header_h, row_h, font_pt = 0.28, 0.22, 12
    if n > 28: row_h, font_pt = 0.20, 11
    if n > 36: row_h, font_pt = 0.18, 10
    tbl.rows[0].height = Inches(header_h)
    for r in range(1, len(tbl.rows)): tbl.rows[r].height = Inches(row_h)

    # 表头 + 数据
    set_header_with_sub(tbl.cell(0,0), "donor", "(chain_residue_atom)")
    set_header_with_sub(tbl.cell(0,1), "acceptor", "(chain_residue_atom)")
    style_cell(tbl.cell(0,2), "distance (Å)", nowrap=True, size_pt=font_pt)
    for row in tbl.rows:
        for cell in row.cells: cell.fill.background()
    for i,(d,a,dist) in enumerate(rows, start=1):
        style_cell(tbl.cell(i,0), d, size_pt=font_pt)
        style_cell(tbl.cell(i,1), a, size_pt=font_pt)
        style_cell(tbl.cell(i,2), dist, size_pt=font_pt, nowrap=True)
def _draw_interface_labels_on_slide(slide, c1, c2, left_labels, right_labels, *,
                                    left_in=5.0, top_in=1.0, width_in=4.8, height_in=6.0,
                                    left_rgb=None, right_rgb=None):
    """在给定 slide 的右侧区域绘制接口残基（每条链最多两列，同页内完成）。"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    AREA_X, AREA_Y, AREA_W, AREA_H = Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    # 转英寸数值便于计算
    area_w_in = width_in

    CHAIN_GAP_IN  = 0.30
    SUBCOL_GAP_IN = 0.20
    FONT_NAME     = "Times New Roman"
    FONT_PT_BASE  = 12; FONT_PT_MIN = 10
    LINE_H_BASE   = 0.26; LINE_H_MIN = 0.18

    BLACK = RGBColor(0,0,0)
    left_color  = RGBColor(*(left_rgb  or (0,0,0)))
    right_color = RGBColor(*(right_rgb or (0,0,0)))

    half_w = (area_w_in - CHAIN_GAP_IN) / 2.0
    def max_rows_per_col(line_h_in):  # 简单估算：按区域高度均分
        return max(1, int(height_in / line_h_in))

    def fit_two_columns(labels):
        font_pt, line_h = FONT_PT_BASE, LINE_H_BASE
        while True:
            m = max_rows_per_col(line_h)
            cols_needed = (len(labels) + m - 1) // m
            if cols_needed <= 2:
                return m, 2 if len(labels) > m else 1, font_pt, line_h
            if line_h > LINE_H_MIN + 1e-6:
                line_h = max(LINE_H_MIN, line_h - 0.02)
            elif font_pt > FONT_PT_MIN:
                font_pt -= 1
            else:
                return m, 2, font_pt, line_h

    def to_columns(labels, m):
        cols=[]
        for i in range(0, min(len(labels), 2*m), m):
            cols.append(labels[i:i+m])
            if len(cols) == 2: break
        return cols

    def draw_chain(x0_in, chain_name, labels, rgb):
        m, ncols, font_pt, line_h = fit_two_columns(labels)
        cols = to_columns(labels, m)
        # 链名
        tb = slide.shapes.add_textbox(Inches(left_in + x0_in), Inches(top_in - 0.30),
                                      Inches(half_w), Inches(0.25))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = f"{chain_name} residues"
        r.font.name = FONT_NAME; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT
        # 列起点
        if ncols == 1:
            xs = [x0_in]
            sub_w = half_w
        else:
            sub_w = (half_w - SUBCOL_GAP_IN) / 2.0
            xs = [x0_in, x0_in + sub_w + SUBCOL_GAP_IN]
        # 逐列渲染
        for idx, col in enumerate(cols):
            x = left_in + xs[idx]
            y = top_in
            for t in col:
                box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(sub_w), Inches(line_h))
                tf = box.text_frame; tf.clear()
                p = tf.paragraphs[0]; rr = p.add_run()
                rr.text = str(t)
                rr.font.name = FONT_NAME; rr.font.size = Pt(font_pt)
                rr.font.color.rgb = (rgb if IFACE_COLOR_TEXT else BLACK)
                y += line_h

    # 两条链：左半区 / 右半区
    draw_chain(0.0,          c1, left_labels,  left_color)
    draw_chain(half_w+CHAIN_GAP_IN, c2, right_labels, right_color)

def add_interface_images_to_ppt(prs, session, chains, svg_dir):
    """
    每个链对：左侧叠加 interface PNG，右侧该链对的文本框（每条链最多两列，同页内完成）。
    """
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    import os

    cmap_by_model, cmap_by_chain = build_chain_color_map(session, getattr(session, "logger", None))
    models = session.models.list(type=AtomicStructure)
    mid = base_mid(models[0].id_string) if models else "#1"

    # 左侧图片区域
    IMG_LEFT_IN  = 0.5
    IMG_TOP_IN   = 2.0
    IMG_W_IN     = 4.8

    for c1, c2 in itertools.combinations(chains, 2):
        # 生成并取 PNG
        _ = generate_interface_images(session, c1, c2, svg_dir,
                                      cmap_by_model=cmap_by_model,
                                      cmap_by_chain=cmap_by_chain,
                                      mid=mid)
        p_cartoon_label = os.path.join(svg_dir, f"interface_{c1}_{c2}_label.png")
        p_atoms_label   = os.path.join(svg_dir, f"interface_{c1}_{c2}_atoms_label.png")
        p_atoms         = os.path.join(svg_dir, f"interface_{c1}_{c2}_atoms.png")
        if not any(os.path.exists(p) for p in (p_cartoon_label, p_atoms_label, p_atoms)):
            continue

        # 右侧文本需要拿到该对的残基标签
        r1_raw, r2_raw = _iface_residues_for_pair(session, c1, c2)
        left_labels  = [_iface_label(nm, num) for nm, num in r1_raw]
        right_labels = [_iface_label(nm, num) for nm, num in r2_raw]

        # 颜色（与链色一致）
        hex1,_ = pick_color(mid, c1, cmap_by_model, cmap_by_chain)
        hex2,_ = pick_color(mid, c2, cmap_by_model, cmap_by_chain)
        left_rgb  = _hex_to_rgb_tuple(hex1)
        right_rgb = _hex_to_rgb_tuple(hex2)

        # 新建一页
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Interface {c1}-{c2} "
        for p in slide.shapes.title.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT

        # 左侧叠加 PNG
        if os.path.exists(p_atoms):
            slide.shapes.add_picture(p_atoms, Inches(IMG_LEFT_IN), Inches(IMG_TOP_IN), width=Inches(IMG_W_IN))
        if os.path.exists(p_atoms_label):
            slide.shapes.add_picture(p_atoms_label, Inches(IMG_LEFT_IN), Inches(IMG_TOP_IN), width=Inches(IMG_W_IN))
        #if os.path.exists(p_cartoon_label):
        #    slide.shapes.add_picture(p_cartoon_label, Inches(IMG_LEFT_IN), Inches(IMG_TOP_IN), width=Inches(IMG_W_IN))

        # 右侧文本（每条链最多两列，单页内完成）
        _draw_interface_labels_on_slide(slide, c1, c2, left_labels, right_labels,
                                        left_in=5.0, top_in=1.0, width_in=4.8, height_in=6.0,
                                        left_rgb=left_rgb, right_rgb=right_rgb)


# =============== 主入口 ===============
def run(session, *args):
    try: cxrun(session, "tool show Log")
    except Exception: pass
    session.logger.info("H_bond.py: run() entered")

    try:
        script_dir = os.path.dirname(__file__)
        out_val = _get_kv(args, "out") or _get_kv(args, "outppt")
        fasta_path = _get_kv(args, "fasta")
        if not fasta_path and os.path.exists(DEFAULT_FASTA):
            fasta_path = DEFAULT_FASTA
            session.logger.info(f"FASTA not provided; using {DEFAULT_FASTA}")
        else:
            session.logger.info(f"FASTA: {fasta_path or 'NONE'}")

        ppt_path, save_txt_path, svg_dir = _derive_out_paths(out_val, script_dir)
        os.makedirs(os.path.dirname(ppt_path), exist_ok=True)
        os.makedirs(os.path.dirname(save_txt_path), exist_ok=True)
        session.logger.info(f"PPT output: {ppt_path}")
        session.logger.info(f"HBond log (saveFile): {save_txt_path}")
        session.logger.info(f"SVG output dir: {svg_dir}")

        use_savefile = not _has_flag(args, "nosavefile")

        chains = sorted({r.chain_id for m in session.models.list(type=AtomicStructure) for r in m.residues})
        if not chains: return

        sel_cmd = 'select ::name="A"::name="C"::name="DA"::name="DC"::name="DG"::name="DT"::name="G"::name="U"'
        session.logger.info("RUN: " + sel_cmd); cxrun(session, sel_cmd)
        has_nucleic = bool(selected_atoms(session))
        session.logger.info("SELECT RESULT: " + ("Nucleic acids selected" if has_nucleic else "Nothing selected"))
        try: cxrun(session, "select clear")
        except Exception: pass

        COMMON = "select true reveal true dashes 6 log true color black radius .05 namingStyle simple"
        if use_savefile:
            if not has_nucleic:
                hb_cmd = f"hbonds intraMol false {COMMON} saveFile {q(save_txt_path)}"
            else:
                hb_cmd = f"hbonds intraMol false {COMMON} saveFile {q(save_txt_path)}"
            session.logger.info("RUN: " + hb_cmd); cxrun(session, hb_cmd)
            log_path = save_txt_path
        else:
            cxrun(session, "log clear")
            if not has_nucleic:
                hb_cmd = f"hbonds intraMol false {COMMON}"
            else:
                hb_cmd = f"hbonds intraMol false {COMMON}"
            session.logger.info("RUN: " + hb_cmd); cxrun(session, hb_cmd)
            tmp_path = os.path.join(tempfile.gettempdir(), "hbonds_log.txt")
            cxrun(session, f'log save {q(tmp_path)}'); session.logger.info(f"LOG saved to {tmp_path}")
            log_path = tmp_path

        rows=[]
        if os.path.exists(log_path):
            with open(log_path, "r", encoding="utf-8", errors="ignore") as f:
                txt = f.read()
            rows = parse_hbond_log(txt)
            if not rows:
                rows = parse_hbond_log_fallback(txt)
        session.logger.info(f"Parsed rows: {len(rows)} from {log_path}")

        from pptx import Presentation
        prs = Presentation()

        add_hbond_all_and_pairs(prs, rows, svg_dir)

        # ====== SALT BRIDGES ======
        try:
            if use_savefile:
                if not has_nucleic:
                    salt_cmd = f"hbonds intraMol false select true saltOnly true reveal true dashes 6 log true color black radius .05 namingStyle simple saveFile {q(save_txt_path.replace('hbond.txt', 'salt.txt'))}"
                else:
                    sel_expr = sel_cmd[len('select '):]
                    salt_cmd = f"hbonds ({sel_expr}) restrict protein select true saltOnly true reveal true dashes 6 log true color black radius .05 namingStyle simple saveFile {q(save_txt_path.replace('hbond.txt', 'salt.txt'))}"
                session.logger.info('RUN: ' + salt_cmd); cxrun(session, salt_cmd)
                salt_log_path = save_txt_path.replace('hbond.txt', 'salt.txt')
            else:
                cxrun(session, "log clear")
                if not has_nucleic:
                    salt_cmd = "hbonds intraMol false select true saltOnly true reveal true dashes 6 log true color black radius .05 namingStyle simple"
                else:
                    sel_expr = sel_cmd[len('select '):]
                    salt_cmd = f"hbonds ({sel_expr}) restrict protein select true saltOnly true reveal true dashes 6 log true color black radius .05 namingStyle simple"
                session.logger.info('RUN: ' + salt_cmd); cxrun(session, salt_cmd)
                tmp_salt = os.path.join(tempfile.gettempdir(), "salts_log.txt")
                cxrun(session, f'log save {q(tmp_salt)}'); session.logger.info(f"SALT LOG saved to {tmp_salt}")
                salt_log_path = tmp_salt

            salt_rows = []
            if os.path.exists(salt_log_path):
                with open(salt_log_path, "r", encoding="utf-8", errors="ignore") as f:
                    salt_txt = f.read()
                salt_rows = parse_hbond_log(salt_txt)
                if not salt_rows:
                    salt_rows = parse_hbond_log_fallback(salt_txt)
            session.logger.info(f"Parsed salt rows: {len(salt_rows)} from {salt_log_path}")

            add_salt_all_and_pairs(prs, salt_rows, svg_dir)
        except Exception as e:
            session.logger.error("SALT stage failed:\n" + "".join(traceback.format_exception(type(e), e, e.__traceback__)))

        # Interface + SVG(+PNG)
        add_interface_pairs(prs, session, fasta_path, svg_dir)

        # =============== 额外：MLP & Coulombic 可视化导出（可选） ===============
        try:
            tmpdir = svg_dir
            mlp_path = os.path.join(tmpdir, "mlp.png")
            coul_path = os.path.join(tmpdir, "coulombic.png")
            mlp_cxs = os.path.join(tmpdir, "mlp.cxs")
            coul_cxs = os.path.join(tmpdir, "coulombic.cxs")

            session.logger.info("RUN: view; ~label; hide atoms;show cartoons; mlp key true")
            cxrun(session, "lighting soft; view #1; ~label; hide atoms;show cartoons;mlp key true;select clear")
            cxrun(session, f"save {q(mlp_path)} width 2000 height 1675 supersample 3 transparentBackground true")
            cxrun(session, f"save {q(mlp_cxs)}")
            session.logger.info(f"MLP image saved: {mlp_path}")

            session.logger.info("RUN: view; coulombic protein key true")
            cxrun(session, "lighting soft; view #1; ~label; hide atoms;show cartoons; coulombic protein key true;select clear")
            cxrun(session, f"save {q(coul_path)} width 2000 height 1675 supersample 3 transparentBackground true")
            cxrun(session, f"save {q(coul_cxs)}")
            session.logger.info(f"Coulombic image saved: {coul_path}")
        except Exception as e:
            session.logger.error(f"MLP/Coulombic export failed: {e}")

        # 把链对图片塞 PPT（可选）
        #add_hbond_salt_images_to_ppt(prs, session, chains, svg_dir, is_salt=False)
        #add_hbond_salt_images_to_ppt(prs, session, chains, svg_dir, is_salt=True)

        add_hbond_salt_images_to_ppt(prs, session, chains, svg_dir, is_salt=False, rows=rows)
        add_hbond_salt_images_to_ppt(prs, session, chains, svg_dir, is_salt=True,  rows=salt_rows)

        add_interface_images_to_ppt(prs, session, chains, svg_dir)

        prs.save(ppt_path)
        session.logger.info(f"PPT saved: {ppt_path}")
        session.logger.info("H_bond.py: done")

    except Exception as e:
        session.logger.error("ERROR in H_bond.py:\n" + "".join(traceback.format_exception(type(e), e, e.__traceback__)))

if 'session' in globals():
    import sys
    try:
        run(session, *sys.argv[1:])
    except Exception as _e:
        try: session.logger.error(f"H_bond.py auto-run failed: {_e}")
        except Exception: pass
