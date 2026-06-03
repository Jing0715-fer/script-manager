#!/usr/bin/env python3
"""
抗原 - 抗体互作序列可视化工具 (E 图格式)
- 空心正方形标记互作残基
- 多结构共享序列，方框竖排排列
- 导出 PNG 和 PPTX 格式（PNG 从 PPTX 转换，保持样式一致）

用法:
    python epitope_visualizer.py --seq seq.fasta --matrix interaction_matrix.xlsx --output output
    # 多结构分析:
    python epitope_visualizer.py --seq s1.fasta,s2.fasta --matrix m1.xlsx,m2.xlsx --labels "Struct1,Struct2" --output combined
    # 仅生成 PPTX（不生成 PNG）:
    python epitope_visualizer.py --seq seq.fasta --matrix matrix.xlsx --output output --no-pptx
"""

import argparse
import re
import os
import subprocess
import shutil
import tempfile
from pathlib import Path
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from lxml import etree


def parse_fasta(fasta_file):
    """解析 FASTA 文件，返回 (名称，序列)"""
    name = ""
    sequence = ""
    with open(fasta_file, 'r') as f:
        for line in f:
            line = line.strip()
            if line.startswith('>'):
                name = line[1:].split()[0] if len(line) > 1 else ""
            else:
                sequence += line
    return name, sequence


def parse_residue(residue_str):
    """解析残基字符串，如 'ASP28(A)' -> (28, 'ASP')"""
    match = re.match(r'([A-Z]{3})(\d+)\([A-Z]\)', residue_str)
    if match:
        return int(match.group(2)), match.group(1)
    return None, None


def parse_chain(residue_str):
    """获取残基所属链"""
    match = re.search(r'\(([A-Z])\)', residue_str)
    return match.group(1) if match else None


def parse_interaction_matrix(matrix_file):
    """解析互作矩阵 Excel 文件，返回：{位置：[抗体链列表]}"""
    df = pd.read_excel(matrix_file, header=None)

    epitope_residues = []
    for i in range(2, 19):
        val = df.iloc[i, 0]
        if pd.notna(val):
            epitope_residues.append(str(val))

    paratope_residues = []
    for j in range(1, 16):
        val = df.iloc[1, j]
        if pd.notna(val):
            paratope_residues.append(str(val))

    epitope_positions = {}

    for i, epitope in enumerate(epitope_residues):
        row_idx = i + 2
        interacting_chains = set()

        for j, paratope in enumerate(paratope_residues):
            col_idx = j + 1
            interaction = df.iloc[row_idx, col_idx]

            if pd.notna(interaction) and str(interaction).strip():
                chain = parse_chain(paratope)
                if chain:
                    interacting_chains.add(chain)

        if interacting_chains:
            pos, _ = parse_residue(epitope)
            if pos is not None:
                epitope_positions[pos] = sorted(list(interacting_chains))

    return epitope_positions


def visualize_epitope(structures, output_file, residues_per_line=50):
    """
    生成抗原 - 抗体互作序列可视化图
    多结构时：序列只显示一次，多个结构的方框竖排排列
    """
    if not structures:
        print("没有结构数据")
        return

    # 找到最长的序列
    max_seq_len = max(len(seq) for _, seq, _, _ in structures)
    num_lines = (max_seq_len + residues_per_line - 1) // residues_per_line

    # 参数设置
    char_spacing = 0.48  # 字符间距
    line_height = 0.65   # 序列行高度
    box_size = 0.18      # 方框尺寸（更小）
    box_spacing = 0.22   # 方框之间的垂直间距

    # 计算高度：需要序列行 + 方框空间（多结构时更多）+ 边距
    num_structs = len(structures)
    # 根据结构数量动态调整行高和总高度
    if num_structs == 1:
        line_height = 0.65
        fig_height = num_lines * line_height + 2.5  # 增加高度用于位置标记
    else:
        line_height = 0.85  # 多结构时增加行高
        fig_height = num_lines * line_height + num_structs * 0.3 + 1.5
    fig_width = min(residues_per_line * char_spacing + 4, 20)

    fig, ax = plt.subplots(figsize=(fig_width, fig_height))
    ax.set_xlim(-1, residues_per_line + 1)
    ax.set_ylim(-num_lines * line_height - 1, 2)  # 上方留更多空间
    ax.axis('off')

    y_seq = 1.0  # 序列行的 Y 坐标（向下移动，给上方标记留空间）

    # 收集所有互作位置及其颜色 {位置：颜色列表}
    all_epitope_colors = {}
    for struct_idx, (name, sequence, epitope_positions, color) in enumerate(structures):
        for pos in epitope_positions.keys():
            if pos not in all_epitope_colors:
                all_epitope_colors[pos] = []
            all_epitope_colors[pos].append(color)

    # 绘制序列行（所有结构共享）
    for line_idx in range(num_lines):
        start = line_idx * residues_per_line
        end = min(start + residues_per_line, max_seq_len)
        y = y_seq - line_idx * line_height

        if start >= max_seq_len:
            break

        # 在序列上方每 10 个位置标记数字（·+ 数字格式，·对齐氨基酸，数字在·上方）
        y_dot = y + 0.45  # ·点的 Y 坐标（靠近序列）
        y_label = y + 0.65  # 数字的 Y 坐标（在·上方）
        for i in range(0, end - start):
            abs_pos = start + i + 1  # 绝对位置（从 1 开始）
            if abs_pos % 10 == 0:  # 每 10 个位置标记
                x_label = i + 0.5
                # 先画·点对齐氨基酸
                ax.text(x_label, y_dot, '·', fontsize=10, ha='center', va='bottom', color='black')
                # 再画数字在·上方
                ax.text(x_label, y_label, str(abs_pos), fontsize=9, ha='center', va='bottom', color='black')

        for i in range(start, end):
            if i >= max_seq_len:
                continue
            pos = i + 1
            aa = sequence[i] if i < len(sequence) else ''
            x = i - start + 0.5

            # 如果该位置是互作残基，使用对应结构的颜色
            # 多个结构共有的位置标记为绿色，单独结构用对应颜色
            if pos in all_epitope_colors:
                if len(all_epitope_colors[pos]) > 1:
                    aa_color = '#00AA00'  # 绿色 - 多个结构共有
                else:
                    aa_color = all_epitope_colors[pos][0]  # 单独结构用对应颜色
            else:
                aa_color = 'black'

            ax.text(x, y, aa, fontsize=10, ha='center', va='center',
                   color=aa_color, weight='bold')

        # 位置标签（行首和行尾）
        ax.text(-0.5, y, str(start + 1), fontsize=8, ha='right', va='center', color='gray')
        ax.text(residues_per_line + 0.3, y, str(end), fontsize=8, ha='left', va='center', color='gray')

    # 计算坐标比例因子，使方框显示为正方形
    # matplotlib 的坐标系中，x 和 y 方向的缩放可能不同，需要补偿
    # x 范围：-1 到 residues_per_line+1 (约 52 单位)
    # y 范围：-num_lines*line_height-0.5 到 1 (约 8 单位)
    # 由于 figure 是宽扁的，x 方向的单位距离比 y 方向短
    # 要让方框看起来是正方形，需要调整宽度
    x_range = residues_per_line + 2
    y_range = num_lines * line_height + 1.5
    # 计算 x 和 y 方向每数据单位的英寸数
    x_scale = fig_width / x_range
    y_scale = fig_height / y_range
    # 宽度缩放因子 = y_scale / x_scale
    box_scale = y_scale / x_scale

    # 绘制每个结构的方框
    for struct_idx, (name, sequence, epitope_positions, color) in enumerate(structures):
        # 计算该结构方框的 Y 偏移（多结构时竖排排列，增加间距让上下排列更清晰）
        y_offset = struct_idx * 0.25  # 增加间距

        for pos in sorted(epitope_positions.keys()):
            # 计算该位置在第几行第几列
            line_idx = (pos - 1) // residues_per_line
            col_in_line = (pos - 1) % residues_per_line

            # 方框的 X 坐标（与上方氨基酸对齐）
            x = col_in_line + 0.5

            # 方框的 Y 坐标（在对应行序列文字的正下方，多结构时竖排）
            seq_line_y = y_seq - line_idx * line_height
            box_y = seq_line_y - 0.45 - y_offset  # 调整方框与序列的间距  # 调整方框与序列的间距

            # 调整方框宽度以补偿坐标比例，使显示为正方形
            rect = patches.Rectangle(
                (x - box_size * box_scale / 2, box_y),
                box_size * box_scale, box_size,
                linewidth=1.8,  # 加粗
                edgecolor=color,
                facecolor='none',
                linestyle='-',
                alpha=1.0
            )
            ax.add_patch(rect)

    plt.tight_layout()
    plt.savefig(output_file, dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()

    total_epitope = sum(len(ep) for _, _, ep, _ in structures)
    print(f"PNG 图像已保存至：{output_file}")
    print(f"结构数：{len(structures)}, 互作残基总数：{total_epitope}")


def pptx_to_png(pptx_file, png_file, dpi=600):
    """使用 PDF 中间转换方法将 PPTX 转换为高质量 PNG"""
    print(f"正在将 PPTX 转换为 PNG (DPI={dpi})...")

    # 检查 LibreOffice 是否可用
    if shutil.which('soffice') is None:
        print(f"警告：LibreOffice 未安装，无法转换 PPTX 为 PNG")
        print(f"PPTX 文件已保存：{pptx_file}")
        return False

    # 创建临时目录
    temp_dir = Path(tempfile.mkdtemp())

    try:
        # 步骤 1: 将 PPTX 转换为 PDF
        pdf_file = temp_dir / f"{Path(pptx_file).stem}.pdf"
        cmd = [
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(temp_dir),
            str(pptx_file)
        ]

        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

        if not pdf_file.exists():
            print(f"PDF 转换失败")
            if result.stderr:
                print(f"LibreOffice 输出：{result.stderr[:200]}")
            return False

        print(f"PDF 已生成：{pdf_file}")

        # 步骤 2: 使用 pdftoppm 将 PDF 转换为 PNG
        if shutil.which('pdftoppm') is not None:
            # 使用 pdftoppm（支持高 DPI 设置）
            output_base = temp_dir / f"{Path(pptx_file).stem}"
            cmd = [
                'pdftoppm',
                '-r', str(dpi),  # 分辨率 DPI
                '-png',          # 输出 PNG 格式
                '-f', '1',       # 从第 1 页开始
                '-l', '1',       # 到第 1 页结束
                str(pdf_file),
                str(output_base)
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

            # 查找生成的 PNG 文件
            generated_png = temp_dir / f"{Path(pptx_file).stem}-1.png"

            if generated_png.exists():
                # 复制到目标位置
                shutil.copy(str(generated_png), str(png_file))
                print(f"PNG 图像已保存至：{png_file}")
                return True
            else:
                print(f"pdftoppm 转换失败，尝试备用方法")
        else:
            print(f"pdftoppm 未安装，使用备用方法")

        # 备用方法：使用 LibreOffice 将 PDF 导出为 PNG
        cmd = [
            'soffice',
            '--headless',
            '--convert-to', 'png',
            '--outdir', str(temp_dir),
            str(pdf_file)
        ]

        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

        # 查找生成的 PNG 文件
        pptx_name = Path(pptx_file).stem
        generated_png = temp_dir / f"{pptx_name}.png"

        if generated_png.exists():
            # 后处理：调整图片大小以达到目标 DPI
            try:
                from PIL import Image

                img = Image.open(str(generated_png))

                # 计算需要缩放的比例（LibreOffice 默认 96 DPI）
                current_dpi = 96
                scale = dpi / current_dpi

                if scale > 1:
                    new_width = int(img.width * scale)
                    new_height = int(img.height * scale)
                    img = img.resize((new_width, new_height), Image.LANCZOS)
                    print(f"已缩放图片至 {new_width}x{new_height} ({dpi} DPI)")

                # 保存处理后的图片
                img.save(str(png_file), 'PNG', dpi=(dpi, dpi))
                print(f"PNG 图像已保存至：{png_file}")

            except ImportError:
                print("PIL 未安装，直接复制图片")
                shutil.copy(str(generated_png), str(png_file))
                print(f"PNG 图像已保存至：{png_file}")

            return True
        else:
            print(f"PNG 转换失败")
            if result.stderr:
                print(f"LibreOffice 输出：{result.stderr[:200]}")
            return False

    except subprocess.TimeoutExpired:
        print("转换超时")
        return False
    except Exception as e:
        print(f"转换出错：{e}")
        return False
    finally:
        # 清理临时文件
        try:
            shutil.rmtree(temp_dir)
        except:
            pass


def create_pptx(structures, output_pptx, residues_per_line=50):
    """创建 PPTX 格式的可视化 - 多结构共享序列，方框竖排"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("错误：需要安装 python-pptx 库：pip install python-pptx")
        return False

    # 计算所需的幻灯片尺寸
    max_seq_len = max(len(seq) for _, seq, _, _ in structures)
    num_lines = (max_seq_len + residues_per_line - 1) // residues_per_line
    num_structs = len(structures)

    # 根据结构数量动态调整行高
    if num_structs == 1:
        line_height = Inches(0.42)
    else:
        line_height = Inches(0.55)

    # 计算内容尺寸
    x_start = Inches(0.8)
    char_width = Inches(0.13)
    content_width = x_start * 2 + residues_per_line * char_width
    content_height = Inches(2.5) + num_lines * line_height  # 顶部序号空间 + 序列行 + 方框空间

    # 创建自定义尺寸的幻灯片（根据内容调整）
    prs = Presentation()
    prs.slide_width = content_width
    prs.slide_height = content_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 参数设置
    y_start = Inches(1.0)  # 向下移动，给上方序号留空间
    box_size = Inches(0.09)
    box_v_spacing = Inches(0.12)  # 方框垂直间距

    # 解析颜色
    def parse_color(color_str):
        """解析颜色字符串，支持十六进制和颜色名称"""
        if color_str.startswith('#'):
            return tuple(int(color_str[i:i+2], 16) for i in (1, 3, 5))
        # 支持常见颜色名称
        color_names = {
            'red': (255, 0, 0),
            'blue': (0, 112, 192),
            'green': (0, 176, 80),
            'orange': (255, 165, 0),
            'purple': (128, 0, 128),
            'black': (0, 0, 0),
            'gray': (128, 128, 128),
            'cyan': (0, 255, 255),
            'magenta': (255, 0, 255),
            'yellow': (255, 255, 0),
        }
        color_lower = color_str.lower().strip()
        if color_lower in color_names:
            return color_names[color_lower]
        # 默认返回红色
        return (204, 0, 0)

    # 收集所有互作位置及其颜色 {位置：颜色列表}
    all_epitope_colors = {}
    for struct_idx, (name, sequence, epitope_positions, color) in enumerate(structures):
        for pos in epitope_positions.keys():
            if pos not in all_epitope_colors:
                all_epitope_colors[pos] = []
            all_epitope_colors[pos].append(color)

    # 绘制序列（所有结构共享）
    max_seq_len = max(len(seq) for _, seq, _, _ in structures)
    num_lines = (max_seq_len + residues_per_line - 1) // residues_per_line

    for line_idx in range(num_lines):
        start = line_idx * residues_per_line
        end = min(start + residues_per_line, max_seq_len)
        y = y_start + line_idx * line_height

        if start >= max_seq_len:
            break

        # 在序列上方每 10 个位置标记数字（·+ 数字格式，·对齐氨基酸，数字在·上方）
        y_dot = y - Inches(0.08)  # ·点的 Y 坐标（靠近序列）
        y_label = y - Inches(0.15)  # 数字的 Y 坐标（在·上方）
        for i in range(0, end - start):
            abs_pos = start + i + 1  # 绝对位置（从 1 开始）
            if abs_pos % 10 == 0:  # 每 10 个位置标记
                x_label = x_start + i * char_width
                # 先画·点对齐氨基酸
                dot_box = slide.shapes.add_textbox(x_label, y_dot, char_width, Inches(0.15))
                tf_dot = dot_box.text_frame
                p_dot = tf_dot.paragraphs[0]
                p_dot.text = '·'
                p_dot.font.size = Pt(8)
                p_dot.font.color.rgb = RGBColor(0, 0, 0)
                p_dot.alignment = PP_ALIGN.CENTER
                # 再画数字在·上方
                label_box = slide.shapes.add_textbox(x_label, y_label, char_width, Inches(0.2))
                tf = label_box.text_frame
                p = tf.paragraphs[0]
                p.text = str(abs_pos)
                p.font.size = Pt(7)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER

        # 位置标签（行首）
        label = slide.shapes.add_textbox(
            x_start - Inches(0.5), y, Inches(0.5), Inches(0.25)
        )
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = str(start + 1)
        p.font.size = Pt(7)
        p.font.color.rgb = RGBColor(128, 128, 128)

        # 序列字母
        for i in range(start, end):
            if i >= max_seq_len:
                continue

            # 获取序列字母
            aa = ''
            for _, sequence, _, _ in structures:
                if i < len(sequence):
                    aa = sequence[i]
                    break

            x = x_start + (i - start) * char_width
            pos = i + 1

            box = slide.shapes.add_textbox(x, y, char_width, Inches(0.25))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = aa
            p.font.size = Pt(8)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # 如果该位置是互作残基，使用对应结构的颜色
            # 多个结构共有的位置标记为绿色，单独结构用对应颜色
            if pos in all_epitope_colors:
                if len(all_epitope_colors[pos]) > 1:
                    aa_color = (0, 170, 0)  # 绿色 - 多个结构共有
                else:
                    aa_color = parse_color(all_epitope_colors[pos][0])
                p.font.color.rgb = RGBColor(*aa_color)
            else:
                p.font.color.rgb = RGBColor(0, 0, 0)

        # 结束位置标签
        end_label = slide.shapes.add_textbox(
            x_start + residues_per_line * char_width, y, Inches(0.4), Inches(0.25)
        )
        tf = end_label.text_frame
        p = tf.paragraphs[0]
        p.text = str(end)
        p.font.size = Pt(7)
        p.font.color.rgb = RGBColor(128, 128, 128)

    # 绘制每个结构的方框 - 在对应氨基酸的正下方，多结构时竖排排列
    for struct_idx, (name, sequence, epitope_positions, color) in enumerate(structures):
        struct_color = parse_color(color)
        # 多结构时，方框竖排排列，紧凑间距
        y_offset = struct_idx * Inches(0.12)

        for pos in sorted(epitope_positions.keys()):
            # 计算该位置在第几行第几列
            line_idx = (pos - 1) // residues_per_line
            col_in_line = (pos - 1) % residues_per_line

            # 方框的 Y 坐标（在对应行序列文字的正下方，多结构时竖排）
            y = y_start + line_idx * line_height + Inches(0.18) + y_offset
            x = x_start + col_in_line * char_width + (char_width - box_size) / 2

            # 使用矩形绘制方框
            shape = slide.shapes.add_shape(
                1,  # msoShapeRectangle
                x,
                y,
                box_size,
                box_size
            )
            # 设置无填充（透明）
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.fill.transparency = 1.0
            # 设置线条颜色
            shape.line.color.rgb = RGBColor(*struct_color)
            shape.line.width = Pt(1.5)

    prs.save(output_pptx)
    print(f"PPTX 文件已保存至：{output_pptx}")

    # 后处理：移除 PPTX 中所有形状的阴影效果
    remove_pptx_shadows(output_pptx)

    return True


def remove_pptx_shadows(pptx_file):
    """移除 PPTX 文件中所有形状的阴影效果，并将白色填充改为无填充"""
    import zipfile
    import shutil
    import re
    from lxml import etree

    temp_dir = Path(tempfile.mkdtemp())

    try:
        # 解压 PPTX 文件
        with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)

        # 查找并修改所有幻灯片
        modified = False
        for slide_file in temp_dir.glob('ppt/slides/slide*.xml'):
            # 读取 XML 内容
            with open(slide_file, 'r', encoding='utf-8') as f:
                content = f.read()

            # 1. 使用正则表达式替换白色填充为无填充
            # 匹配 <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            old_content = content
            content = re.sub(
                r'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>',
                r'<a:noFill/>',
                content
            )
            if content != old_content:
                modified = True
                print(f"已替换白色填充为无填充")

            # 2. 移除所有 outerShdw 阴影效果
            old_content = content
            content = re.sub(
                r'<a:outerShdw[^>]*/>|<a:outerShdw[^>]*></a:outerShdw>',
                '',
                content
            )
            if content != old_content:
                modified = True
                print(f"已移除 outerShdw 阴影效果")

            # 3. 移除形状的 style 元素（包含 effectRef 阴影引用）
            # p:style 中的 a:effectRef 会导致 PowerPoint 渲染阴影
            old_content = content
            content = re.sub(
                r'<p:style>.*?</p:style>',
                '',
                content,
                flags=re.DOTALL
            )
            if content != old_content:
                modified = True
                print(f"已移除形状样式（包含 effectRef 阴影引用）")

            # 保存修改后的 XML
            with open(slide_file, 'w', encoding='utf-8') as f:
                f.write(content)

        # 重新打包 PPTX
        backup_file = str(pptx_file) + '.bak'
        shutil.move(pptx_file, backup_file)

        with zipfile.ZipFile(pptx_file, 'w') as zip_ref:
            for file_path in temp_dir.rglob('*'):
                if file_path.is_file():
                    arc_name = file_path.relative_to(temp_dir)
                    zip_ref.write(file_path, arc_name)

        # 删除备份
        os.remove(backup_file)

        if modified:
            print(f"已移除 PPTX 中的阴影效果")

    except Exception as e:
        print(f"移除阴影失败：{e}")
        # 如果失败，尝试恢复备份
        backup_file = str(pptx_file) + '.bak'
        if os.path.exists(backup_file):
            shutil.move(backup_file, pptx_file)
    finally:
        # 清理临时文件
        try:
            shutil.rmtree(temp_dir)
        except:
            pass


def load_structure_data(seq_file, matrix_file, color=None):
    """加载单个结构的数据"""
    name, sequence = parse_fasta(seq_file)
    epitope_positions = parse_interaction_matrix(matrix_file)
    if not name:
        name = os.path.splitext(os.path.basename(seq_file))[0]
    return name, sequence, epitope_positions


def main():
    parser = argparse.ArgumentParser(
        description='抗原 - 抗体互作序列可视化工具 (E 图格式)',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  单结构:
    python epitope_visualizer.py --seq seq.fasta --matrix matrix.xlsx

  多结构分析:
    python epitope_visualizer.py --seq s1.fasta,s2.fasta --matrix m1.xlsx,m2.xlsx --labels "Struct1,Struct2"
    python epitope_visualizer.py --seq s1.fasta,s2.fasta --matrix m1.xlsx,m2.xlsx --labels "S1,S2" --colors "red,blue"
        """
    )
    parser.add_argument('--seq', required=True,
                       help='输入 FASTA 序列文件（多结构用逗号分隔）')
    parser.add_argument('--matrix', required=True,
                       help='输入互作矩阵 Excel 文件（多结构用逗号分隔）')
    parser.add_argument('--output', default='epitope_map',
                       help='输出文件路径前缀 (默认：epitope_map)')
    parser.add_argument('--labels', default='',
                       help='结构标签，用逗号分隔（可选）')
    parser.add_argument('--colors', default='',
                       help='每个结构的标记颜色，用逗号分隔（可选）')
    parser.add_argument('--residues-per-line', type=int, default=50,
                       help='每行显示的残基数 (默认：50)')
    parser.add_argument('--no-pptx', action='store_true',
                       help='不生成 PPTX 文件')

    args = parser.parse_args()

    seq_files = [s.strip() for s in args.seq.split(',')]
    matrix_files = [m.strip() for m in args.matrix.split(',')]
    labels = [l.strip() for l in args.labels.split(',')] if args.labels else []
    colors = [c.strip() for c in args.colors.split(',')] if args.colors else None

    if len(seq_files) != len(matrix_files):
        print("错误：序列文件数量和矩阵文件数量必须一致")
        return 1

    structures = []
    for i, (seq_file, matrix_file) in enumerate(zip(seq_files, matrix_files)):
        label = labels[i] if i < len(labels) else None
        color = colors[i] if colors and i < len(colors) else None
        name, sequence, epitope_positions = load_structure_data(seq_file, matrix_file, color)
        if label:
            name = label
        if not color:
            color = None
        structures.append((name, sequence, epitope_positions, color))
        print(f"结构 {i+1}: {name}")
        print(f"  序列长度：{len(sequence)}")
        print(f"  互作残基：{len(epitope_positions)}")

    # 默认颜色
    default_colors = ['#CC0000', '#0066CC', '#009900', '#FF6600', '#8B008B']
    if colors is None:
        colors = default_colors

    # 确保每个结构都有颜色
    final_structures = []
    for i, (n, s, e, c) in enumerate(structures):
        if c is None:
            c = colors[i % len(colors)]
        final_structures.append((n, s, e, c))
    structures = final_structures

    print(f"\n生成可视化图像...")

    # 输出文件路径
    output_base = args.output
    if not output_base.endswith('.pptx'):
        pptx_file = output_base + '.pptx'
        png_file = output_base + '.png'
    else:
        pptx_file = output_base
        png_file = output_base.replace('.pptx', '.png')

    # 生成 PPTX
    if not args.no_pptx:
        create_pptx(
            structures,
            pptx_file,
            residues_per_line=args.residues_per_line
        )
        # 将 PPTX 转换为 PNG（使用 600 DPI 高质量）
        pptx_to_png(pptx_file, png_file, dpi=600)
    else:
        # 仅生成 PNG（使用 matplotlib）
        visualize_epitope(
            structures,
            png_file,
            residues_per_line=args.residues_per_line
        )

    return 0


if __name__ == '__main__':
    exit(main() or 0)
